import os
import re
import shutil
import zipfile
from pptx import Presentation
from pptx.text.text import Font
from pathlib import Path
from django.db import models
from django.core.exceptions import ValidationError
from django.utils.translation import gettext_lazy as _
from django.utils import timezone
from django.utils.html import format_html, format_html_join


BASE_DIR = Path(__file__).resolve().parent.parent
MEDIA_ROOT  = BASE_DIR / 'media'


def find_2rd_item_in_maps(item, maps):
    for item1, item2 in maps:
        if item1 == item: 
            return item2

def validate_pptx(value):
    extension = value.name.split('.')[-1]
    if extension != 'pptx':
        raise ValidationError(
            _(value.name+'不是ppt文件(.pptx)'),
            params={'value': value.name},
        )

## text_frame: family, size, bold, italic
def get_rPr_font(text_frame):
    txBody = text_frame._txBody
    # print(len(txBody.p_lst))
    for p in txBody.p_lst:
        for elm in p.content_children:
            return Font(elm.get_or_add_rPr())


Standard_Color_Choices = [
    ('C00000','深红'),
    ('FF0000','红色'),
    ('FFC000','橙色'),
    ('FFFF00','黄色'),
    ('92D050','浅绿'),
    ('00B050','绿色'),
    ('00B0F0','浅蓝'),
    ('0070C0','蓝色'),
    ('002060','深蓝'),
    ('7030A0','紫色'),
]

LINE_NUM_CHOICES = [
    ('2','2'),
    ('3','3'),
    ('4','4'),
    ('5','5'),
    ('6','6'),
]

Slide_Layout_Names = [
    ("标题幻灯片", "标题幻灯片"),
    ("标题和内容", "标题和内容"),
    ("节标题", "节标题"),
    ("两栏内容", "两栏内容"),
    ("比较", "比较"),
    ("仅标题", "仅标题"),
    ("内容与标题", "内容与标题"),
    ("图片与标题", "图片与标题"),
    # ("空白", "空白"),
    # ("标题和竖排文字", "标题和竖排文字"),
    # ("竖排标题与文本", "竖排标题与文本"),
]

Slide_Names = [
    ('0','第一页幻灯片'),
    ('1','第二页幻灯片'),
    ('2','第三页幻灯片'),
    ('3','第四页幻灯片'),
]

Background_Types = [
    ('SOLID (1)','纯色填充'),
    ('PATTERNED (2)','图案填充'),
]

Shape_Names = [
    ('标题区域','标题区域'),
    ('副标题区域','副标题区域'),
    ('内容区域','内容区域'),
    ('左侧内容区域','左侧内容区域'),
    ('右侧内容区域','右侧内容区域'),
]


FONT_NAME_CHOICES = [
    ('黑体', '黑体'),
    ('华文仿宋', '华文仿宋'),
    ('楷体', '楷体'),
    ('幼圆', '幼圆'),
    ('微软雅黑', '微软雅黑'),
    ('隶书', '隶书'),
]


FONT_SIZE_CHOICES = [
    ('26','26' ),
    ('36','36' ),
    ('48','48' ),
    ('72','72' ),
]


def clean_empty_item(item, errors_dict, target_key):
    if item == '': errors_dict[target_key] = _('不能为空')

####################### PPT 题目 #########################
class PPTQuestion(models.Model):

    def __str__(self):
        return 'PPT操作' + str(self.id) 

    class Meta:
        db_table = 'ppt_table'  
        verbose_name_plural = '题目-PPT操作'  
        verbose_name = '题目-PPT操作' 

    def base_path_(self):
        return os.path.join(MEDIA_ROOT, 'upload_pptx', str(self.id))

    def zip_path_(self):
        return os.path.join(MEDIA_ROOT, 'upload_pptx',str(self.id),'ppt操作题目.zip')

    def save(self, *args, **kwargs):
        super().save(*args, **kwargs)  
        if not os.path.exists(self.base_path_()):  os.mkdir(self.base_path_())

        zip_path = self.zip_path_()
        if zip_path:
            if os.path.exists(zip_path): os.remove(zip_path)
            zf = zipfile.ZipFile(zip_path, 'w')
            zf.write(self.upload_pptx.path,'{}/ppt.pptx'.format('ppt操作题目'))
            # if self.image_op:
                # zf.write(self.upload_image_file.path,'{}/{}'.format('ppt操作题目', self.upload_image_file.name.split('/')[-1]))
            zf.close()

    # 题目内容
    def question_content(self):
        pre_description = '打开下载的ppt.pptx,执行以下操作:\n'
        result_list = []
        if self.slide_layout_op:
            result_list.append(find_2rd_item_in_maps(self.slide_layout_target_slide,Slide_Names)+'的版式设置为"'+\
                self.slide_layout_name+'".')

        if self.text_op:
            result_list.append(find_2rd_item_in_maps(self.text_target_slide,Slide_Names)+\
                self.text_target_shape_1+'输入"'+self.text_target_content_1+'", '+\
                self.text_target_shape_2+'输入"'+self.text_target_content_2+'"'+\
                '.')

        if self.font_op:
            additions = ''
            if self.font_bold: additions = additions + '、粗体'
            if self.font_italic: additions = additions + '、斜体'
            result_list.append('设置'+find_2rd_item_in_maps(self.font_target_slide,Slide_Names)+\
                '"'+self.font_target_content+'"内容的字体'+self.font_name+'、'+\
                '字号'+self.font_size+'、'+ \
                '颜色标准色'+ find_2rd_item_in_maps(self.font_color, Standard_Color_Choices) +\
                additions + '.')

        if self.slide_background_op:
            if self.slide_background_type == 'PATTERNED (2)':
                result_list.append(find_2rd_item_in_maps(self.slide_background_slide,Slide_Names)+\
                    '的背景格式设为图案填充，前景标准色'+\
                find_2rd_item_in_maps(self.slide_background_fore, Standard_Color_Choices)+\
                    '，背景标准色'+\
                find_2rd_item_in_maps(self.slide_background_back, Standard_Color_Choices)+\
                    '.')
            elif self.slide_background_type != 'PATTERNED (2)':
                result_list.append(find_2rd_item_in_maps(self.slide_background_slide,Slide_Names)+\
                    '的背景格式设为纯色填充，标准色'+\
                find_2rd_item_in_maps(self.slide_background_fore, Standard_Color_Choices)+\
                    '.')

        if self.notes_slide_op:
            result_list.append(find_2rd_item_in_maps(self.notes_slide_target_slide,Slide_Names)+\
                '添加备注，内容是"'+self.notes_slide_content+\
                '".')

        html_table = format_html("")
        if self.table_op:
            result_list.append(find_2rd_item_in_maps(self.table_target_slide,Slide_Names)+\
                self.table_target_shape+'中插入'+self.table_rows+'行'+self.table_columns+'列表格，'+\
                    '具体内容是如下：')
            html_table = self.html_table_()

        return format_html(pre_description) + \
                format_html("<ol>") + \
                format_html_join(
                '\n', '<li style="color:black;">{}</li>',
                ((x,) for x in result_list)
                ) \
                + format_html("</ol>") + html_table
    question_content.short_description = '题目内容'

    def html_table_(self):
        html_table = format_html("<div><table border='2' align='center'>")
        for rows in self.table_content.split('\n'):
            html_table = html_table+ format_html("<tr >")
            for item in rows.split(','):
                html_table = html_table+ format_html("<td>&nbsp"+item+"&nbsp</td>")
            html_table = html_table+ format_html("</tr>")
        html_table = html_table+ format_html("</table></div>")
        return html_table

    def file_path_(self):
        if self.upload_pptx:
            answer_file_path = self.upload_pptx.path
        else:
            answer_file_path = ''
        return answer_file_path
    file_path_.short_description = '文件地址'


    def score_(self, answer_file_path):
        result_list = []
        ## check the files
        if answer_file_path and os.path.exists(answer_file_path):
            try:
                prs = Presentation(answer_file_path)
                slides_len = len(prs.slides)
            except:
                return ['文件打开异常'], 0

            if self.slide_layout_op:
                slide_id = int(self.slide_layout_target_slide)
                result = 'slide_layout::'+str(slide_id)+'::'
                if slide_id+1 <= slides_len:
                    if self.slide_layout_name == prs.slides[slide_id].slide_layout.name:
                        result = result + self.slide_layout_name
                result_list.append(result)


            if self.text_op:
                slide_id = int(self.text_target_slide)
                result = 'text::'+str(slide_id)+'::'
                contents = [self.text_target_content_1, self.text_target_content_2]
                if slide_id+1 <= slides_len:
                    for shape in prs.slides[slide_id].shapes:
                        if shape.has_text_frame:
                            if (shape.text_frame.text in contents):
                                result = result + shape.text_frame.text + '+'
                result_list.append(result)


            if self.font_op:
                slide_id = int(self.font_target_slide)
                result = 'font::'+str(slide_id)+'::'
                if slide_id+1 <= slides_len:
                     for shape in prs.slides[slide_id].shapes:
                        if shape.has_text_frame and (shape.text_frame.text==self.font_target_content):
                            f = get_rPr_font(shape.text_frame)
                            name, size, bold, color = f.name, f.size, f.bold, f.color
                            # print(name, self.font_name)
                            # print(str(int(size.pt)), self.font_size)
                            if name == self.font_name:
                                result = result + name + '+'
                            if size and str(int(size.pt)) == self.font_size:
                                result = result + self.font_size + '+'
                            if bold and self.font_bold:
                                result = result + '粗体' + '+'
                            if str(color.type) == 'RGB (1)' and str(color.rgb)==self.font_color:
                                result = result + self.font_color
                            break
                result_list.append(result)


            if self.slide_background_op:
                slide_id = int(self.slide_background_slide)
                result = 'slide_background::'+str(slide_id)+'::'
                if slide_id+1 <= slides_len:
                    fill = prs.slides[slide_id].background.fill
                    if self.slide_background_type == str(fill.type):
                        result = result + self.slide_background_type + '+'
                        if str(fill.type) == 'PATTERNED (2)':
                            if str(fill.fore_color.type) == 'RGB (1)' and  \
                                self.slide_background_fore == str(fill.fore_color.rgb):
                                result = result + self.slide_background_fore + '+'
                            if str(fill.back_color.type) == 'RGB (1)' and  \
                                self.slide_background_back == str(fill.back_color.rgb):
                                result = result + self.slide_background_back
                        elif str(fill.type) == 'SOLID (1)':
                            if str(fill.fore_color.type) == 'RGB (1)' and  \
                                self.slide_background_fore == str(fill.fore_color.rgb):
                                result = result + self.slide_background_fore
                result_list.append(result)


            if self.notes_slide_op:
                slide_id = int(self.notes_slide_target_slide)
                result = 'notes_slide::'+str(slide_id)+'::'
                if slide_id+1 <= slides_len:
                    if self.notes_slide_content == prs.slides[slide_id].notes_slide.notes_text_frame.text:
                        result = result + self.notes_slide_content
                result_list.append(result)


            if self.table_op:
                slide_id = int(self.table_target_slide)
                result = 'table::'+str(slide_id)+'::'
                if slide_id+1 <= slides_len:
                    for shape in prs.slides[slide_id].shapes:
                        if shape.has_table:
                            result = result + str(len(shape.table.columns))+'+'+\
                                str(len(shape.table.rows))+'+'+shape.table.cell(0,0).text
                result_list.append(result)

        else:
            result_list.append('Nothing to score')

        score = len([x for x in result_list if x.split('::')[1]])*1.0/len(result_list)
        return result_list, score


    def test_(self):
        ## only for development
        if self.upload_pptx:
            answer_file_path = self.upload_pptx.path
        else:
            answer_file_path = ''

        result_list, score = self.score_(answer_file_path)
        return format_html("<ol>") + \
                format_html_join(
                '\n', '<li style="color:black;">{}</li>',
                ((x,) for x in result_list)
                ) \
                + format_html("</ol>")
    test_.short_description = '测试结果'


    # def save(self, *args, **kwargs):
    #     super().save(*args, **kwargs)  
        

    def clean(self):
        error_dict = {}

        if self.slide_layout_op:
            clean_empty_item(self.slide_layout_target_slide, error_dict, 'slide_layout_target_slide')
            clean_empty_item(self.slide_layout_name, error_dict, 'slide_layout_name')

        if self.text_op:
            clean_empty_item(self.text_target_slide, error_dict, 'text_target_slide')
            clean_empty_item(self.text_target_shape_1, error_dict, 'text_target_shape_1')
            clean_empty_item(self.text_target_content_1, error_dict, 'text_target_content_1')
            clean_empty_item(self.text_target_shape_2, error_dict, 'text_target_shape_2')
            clean_empty_item(self.text_target_content_2, error_dict, 'text_target_content_2')
            if self.text_target_shape_1 == self.text_target_shape_2:
                error_dict['text_target_shape_1'] = _('两个区域不能相同')
                error_dict['text_target_shape_2'] = _('两个区域不能相同')


        if self.font_op:
            clean_empty_item(self.font_target_slide, error_dict, 'font_target_slide')
            clean_empty_item(self.font_target_shape, error_dict, 'font_target_shape')
            clean_empty_item(self.font_target_content, error_dict, 'font_target_content')
            clean_empty_item(self.font_name, error_dict, 'font_name')
            clean_empty_item(self.font_size, error_dict, 'font_size')
            clean_empty_item(self.font_color, error_dict, 'font_color')

        if self.slide_background_op:
            clean_empty_item(self.slide_background_slide, error_dict, 'slide_background_slide')
            clean_empty_item(self.slide_background_type, error_dict, 'slide_background_type')
            clean_empty_item(self.slide_background_fore, error_dict, 'slide_background_fore')
            clean_empty_item(self.slide_background_back, error_dict, 'slide_background_back')
            if self.slide_background_fore == self.slide_background_back:
                error_dict['slide_background_fore'] = _('前景和背景颜色不能相同')
                error_dict['slide_background_back'] = _('前景和背景颜色不能相同')


        if self.notes_slide_op:
            clean_empty_item(self.notes_slide_target_slide, error_dict, 'notes_slide_target_slide')
            clean_empty_item(self.notes_slide_content, error_dict, 'notes_slide_content')

        if self.table_op:
            clean_empty_item(self.table_target_slide, error_dict, 'table_target_slide')
            clean_empty_item(self.table_target_shape, error_dict, 'table_target_shape')
            clean_empty_item(self.notes_slide_content, error_dict, 'notes_slide_content')
            clean_empty_item(self.table_rows, error_dict, 'table_rows')
            clean_empty_item(self.table_columns, error_dict, 'table_columns')
            clean_empty_item(self.table_content, error_dict, 'table_content')
            contents = self.table_content.split('\n')
            if len(contents) != int(self.table_rows):
                error_dict['table_content'] = _('表格内容和行数不相等')
            for c in contents:
                c = [x for x in c.split(',') if x]
                if len(c) == int(self.table_columns): continue
                error_dict['table_content'] = _('表格内容和列数不相等')
                break

        if len(error_dict)>0:
            raise ValidationError(error_dict)


    #########
    pub_date = models.DateTimeField('创建时间', default=timezone.now)

    ## 
    upload_pptx = models.FileField(verbose_name='上传PPT文件[.pptx]', upload_to='upload_pptx/', 
         validators=[validate_pptx], null=True)

    # 选择版式
    # 题目示例：第X页幻灯片的版式改为“XXXX”。
    slide_layout_op = models.BooleanField('考查版式设置？', default=False)
    slide_layout_target_slide = models.CharField('要操作的幻灯片', choices=Slide_Names ,max_length=50, blank=True, default='0')
    slide_layout_name = models.CharField('版式名称', choices=Slide_Layout_Names ,max_length=50, blank=True, default='标题幻灯片')


    # 修改文字内容
    text_op = models.BooleanField('考查修改文字内容？', default=False)
    text_target_slide = models.CharField('要操作的幻灯片', choices=Slide_Names ,max_length=50, blank=True, default='0')
    text_target_shape_1 = models.CharField('区域1', choices=Shape_Names ,max_length=50, blank=True, default='标题')
    text_target_content_1 = models.TextField('区域1内容',  blank=True, default='大标题')
    text_target_shape_2 = models.CharField('区域2', choices=Shape_Names ,max_length=50, blank=True, default='内容')
    text_target_content_2 = models.TextField('区域2内容',  blank=True, default='小标题')


    # 修改字体格式
    font_op = models.BooleanField('考查修改字体格式？', default=False)
    font_target_slide = models.CharField('要操作的幻灯片', choices=Slide_Names ,max_length=50, blank=True, default='0')
    font_target_shape = models.CharField('文字区域', choices=Shape_Names ,max_length=50, blank=True, default='标题')
    font_target_content = models.TextField('文字内容', blank=True, default='')
    font_name = models.CharField('字体', choices=FONT_NAME_CHOICES ,max_length=50, blank=True, )
    font_size = models.CharField('字号[磅]', choices=FONT_SIZE_CHOICES ,max_length=50, blank=True, )
    font_color = models.CharField('颜色', choices=Standard_Color_Choices ,max_length=50, blank=True,default='红色')
    font_bold = models.BooleanField('是否粗体？', default=False)
    font_italic = models.BooleanField('是否斜体？', default=False)


    # 设置幻灯片背景格式
    # 题目示例：第X页幻灯片设置背景格式为"XX",标准色“XXXX”。
    slide_background_op = models.BooleanField('考查幻灯片背景格式？', default=False)
    slide_background_slide = models.CharField('要操作的幻灯片', choices=Slide_Names, max_length=50, blank=True, default='1')
    slide_background_type = models.CharField('背景类型',  choices=Background_Types, max_length=50, blank=True, default='PATTERNED (2)')
    slide_background_fore = models.CharField('前景颜色',  choices=Standard_Color_Choices, max_length=50, blank=True,default='黄色')
    slide_background_back = models.CharField('背景颜色',  choices=Standard_Color_Choices, max_length=50, blank=True,default='浅蓝')


    # 添加备注
    # 题目示例：第X页幻灯片添加备注，内容为“XXXX”。
    notes_slide_op = models.BooleanField('考查添加备注页？', default=False)
    notes_slide_target_slide = models.CharField('要操作的幻灯片', choices=Slide_Names ,max_length=50, blank=True, default='1')
    notes_slide_content = models.TextField('备注内容',  blank=True, default='备注内容1')

    # 插入表格及内容
    table_op = models.BooleanField('考查插入表格？', default=False)
    table_target_slide = models.CharField('要操作的幻灯片', choices=Slide_Names ,max_length=50, blank=True, default='1')
    table_target_shape = models.CharField('表格插入区域', choices=Shape_Names ,max_length=50, blank=True, default='左侧内容')
    table_rows = models.CharField('表格行数',max_length=50, choices=LINE_NUM_CHOICES, blank=True, default='2')
    table_columns = models.CharField('表格列数',max_length=50, choices=LINE_NUM_CHOICES, blank=True, default='4')
    table_content = models.TextField('表格内容[用逗号和换行分割]',  blank=True, default='橘子,苹果,香蕉,桃子\n黄色,红色,黄色,红色')


